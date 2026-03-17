import streamlit as st
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
import io
from datetime import datetime, timedelta
from supabase import create_client, Client

# --- CONFIGURAÇÕES DA PÁGINA ---
# Definimos o título, ícone e o layout para garantir que o sistema ocupe bem o espaço da tela.
st.set_page_config(
    page_title="Ecoa: Escuta Organizacional",
    page_icon="🌱",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- ESTILIZAÇÃO CSS AVANÇADA (LAYOUT MODERNO E PREMIUM) ---
# Aqui injetamos CSS para transformar a interface padrão em um sistema corporativo elegante.
st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        html, body, [class*="css"] {
            font-family: 'Inter', sans-serif;
            color: #1e293b;
        }

        /* Ocultar elementos padrões do Streamlit */
        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}
        header {visibility: hidden;}
        
        /* Ajuste do container principal */
        .block-container {
            padding-top: 2rem;
            padding-bottom: 2rem;
            max-width: 1200px;
        }

        /* Estilização da Sidebar (Menu Lateral) */
        [data-testid="stSidebar"] {
            background-color: #f8fafc;
            border-right: 1px solid #e2e8f0;
        }
        
        /* Estilização dos itens do Menu (Radio) */
        div.stRadio > div {
            gap: 8px;
        }
        
        div.stRadio label {
            background-color: transparent;
            border-radius: 8px;
            padding: 10px 15px;
            transition: all 0.2s ease;
            border: 1px solid transparent;
            cursor: pointer;
            width: 100%;
        }

        div.stRadio label:hover {
            background-color: #f1f5f9;
            border-color: #e2e8f0;
        }

        div.stRadio label[data-selected="true"] {
            background-color: #eff6ff !important;
            border-color: #bfdbfe !important;
            color: #1e40af !important;
            font-weight: 600;
        }

        /* Botões customizados */
        .stButton>button {
            border-radius: 8px;
            font-weight: 600;
            padding: 0.6rem 1.2rem;
            transition: all 0.3s;
            border: 1px solid #cbd5e1;
        }
        
        .stButton>button:hover {
            border-color: #3b82f6;
            background-color: #f0f9ff;
            color: #3b82f6;
        }

        /* Cards de Módulo */
        .module-card {
            background-color: #ffffff;
            padding: 2rem;
            border-radius: 12px;
            border: 1px solid #e2e8f0;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
            margin-bottom: 1.5rem;
        }

        /* Login Screen Styles */
        .login-box {
            background-color: white;
            padding: 3rem;
            border-radius: 20px;
            box-shadow: 0 10px 25px -5px rgba(0, 0, 0, 0.1);
            border: 1px solid #f1f5f9;
        }
    </style>
""", unsafe_allow_html=True)

# --- CONEXÃO AUTOMÁTICA COM O SUPABASE (SECRETS) ---
@st.cache_resource
def init_connection():
    try:
        url = st.secrets["SUPABASE_URL"]
        key = st.secrets["SUPABASE_KEY"]
        return create_client(url, key)
    except Exception:
        return None

supabase = init_connection()

# --- INICIALIZAÇÃO DE ESTADOS DO SISTEMA ---
if 'autenticado' not in st.session_state:
    st.session_state.autenticado = False

if 'usuario_logado' not in st.session_state:
    st.session_state.usuario_logado = None

if 'perfil_logado' not in st.session_state:
    st.session_state.perfil_logado = None

# Base de usuários inicial (Será sincronizada com o Supabase futuramente)
if 'usuarios' not in st.session_state:
    st.session_state.usuarios = pd.DataFrame({
        "Nome": ["Cris Lima", "Juliana", "Analista RH"],
        "E-mail": ["cris@ecoa.app", "ju@ecoa.app", "rh@empresa.com"],
        "Senha": ["admin123", "senha456", "cliente789"],
        "Perfil": ["Admin Master", "Analista RH", "Cliente (Leitura)"],
        "Qtd. Links": [10000, 2000, 500],
        "Status": ["Ativo", "Ativo", "Ativo"]
    })

# Perguntas baseadas em metodologias de mercado (Gallup Q12, GPTW, Aristotle)
if 'perguntas' not in st.session_state:
    st.session_state.perguntas = pd.DataFrame({
        "Ativa": [True] * 19,
        "Pilar Estratégico": [
            "Liderança & Gestão", "Liderança & Gestão", "Liderança & Gestão",
            "Segurança Psicológica", "Segurança Psicológica",
            "Bem-estar & Saúde Mental", "Bem-estar & Saúde Mental",
            "Reconhecimento & Valorização", "Reconhecimento & Valorização",
            "Desenvolvimento & Crescimento", "Desenvolvimento & Crescimento",
            "Comunicação & Transparência", "Comunicação & Transparência",
            "Diversidade & Inclusão", "Diversidade & Inclusão",
            "Propósito & Orgulho", "Propósito & Orgulho",
            "Infraestrutura & Ferramentas", "eNPS"
        ],
        "Texto da Pergunta": [
            "Meu gestor direto estabelece expectativas claras sobre o que é esperado do meu trabalho.",
            "Recebo feedbacks construtivos e regulares que me ajudam a melhorar meu desempenho.",
            "Sinto que meu gestor se importa comigo como pessoa, não apenas como um número.",
            "Nesta equipe, me sinto seguro para assumir riscos sem medo de punição ou retaliação.",
            "Sinto-me à vontade para expressar opiniões divergentes ou trazer problemas difíceis à tona.",
            "A organização incentiva o equilíbrio saudável entre minha vida pessoal e profissional.",
            "O volume de trabalho exigido me permite realizar minhas tarefas sem comprometer minha saúde mental.",
            "Nos últimos meses, recebi reconhecimento ou elogios genuínos por um trabalho bem feito.",
            "Acredito que as promoções e méritos nesta empresa são justos e baseados em mérito.",
            "Tive oportunidades claras de aprender e crescer profissionalmente no último ano.",
            "Meu trabalho atual desafia minhas habilidades e me permite usar meu pleno potencial.",
            "A liderança mantém os colaboradores informados sobre decisões estratégicas importantes.",
            "Sinto que há um canal aberto e seguro para que minha voz chegue até a diretoria.",
            "As pessoas aqui são tratadas com respeito, independente de raça, gênero ou idade.",
            "Sinto que posso ser eu mesmo(a) no ambiente de trabalho sem precisar de 'máscaras'.",
            "O propósito da empresa me faz sentir que meu trabalho diário tem um significado real.",
            "Tenho muito orgulho de dizer a outras pessoas que trabalho nesta organização.",
            "Tenho as ferramentas e o ambiente necessário para realizar meu trabalho com excelência.",
            "De 0 a 10, o quanto recomendaria nossa empresa como um bom lugar para trabalhar?"
        ]
    })

if 'mensagem_padrao' not in st.session_state:
    st.session_state.mensagem_padrao = "Olá, {nome}! Sua voz é fundamental. Responda à nossa pesquisa até {data_validade}.\nAcesse: {link_pesquisa}"

if 'empresa_atual' not in st.session_state:
    st.session_state.empresa_atual = "Sua Empresa"

if 'data_validade' not in st.session_state:
    st.session_state.data_validade = datetime.today().date() + timedelta(days=15)

if 'logo_personalizada' not in st.session_state:
    st.session_state.logo_personalizada = None

# Banco de dados de respostas simulado para o Dashboard
if 'dados_historicos' not in st.session_state:
    st.session_state.dados_historicos = pd.DataFrame({
        "departamento": ["RH", "RH", "TI", "TI", "Comercial", "Comercial", "Financeiro", "Logística", "Operações"],
        "lideranca": [4.5, 4.8, 3.2, 3.5, 4.7, 4.6, 4.0, 3.8, 3.9],
        "reconhecimento": [4.0, 4.2, 2.8, 3.1, 4.5, 4.3, 3.7, 3.5, 3.6],
        "comunicacao": [4.2, 4.4, 3.0, 3.3, 4.6, 4.5, 3.8, 3.6, 3.7],
        "enps": [10, 10, 6, 7, 10, 9, 8, 7, 8]
    })

# =====================================================================
# TELA DE LOGIN (SECURITY FIRST)
# =====================================================================
if not st.session_state.autenticado:
    _, col_login, _ = st.columns([1, 1.2, 1])
    
    with col_login:
        st.write("")
        st.write("")
        st.write("")
        st.markdown("<h1 style='text-align: center; color: #1e3a8a; font-size: 3rem;'>🌱 Ecoa</h1>", unsafe_allow_html=True)
        st.markdown("<p style='text-align: center; color: #64748b; margin-top: -15px;'>Plataforma de Escuta Organizacional</p>", unsafe_allow_html=True)
        
        with st.form("login_form"):
            st.markdown("#### Login do Usuário")
            user_email = st.text_input("E-mail", placeholder="seu@email.com")
            user_pass = st.text_input("Senha", type="password", placeholder="••••••••")
            
            entrar = st.form_submit_button("Acessar Plataforma", use_container_width=True)
            
            if entrar:
                u_base = st.session_state.usuarios
                valido = u_base[(u_base['E-mail'] == user_email) & (u_base['Senha'] == user_pass) & (u_base['Status'] == 'Ativo')]
                
                if not valido.empty:
                    st.session_state.autenticado = True
                    st.session_state.usuario_logado = valido.iloc[0]['Nome']
                    st.session_state.perfil_logado = valido.iloc[0]['Perfil']
                    st.rerun()
                else:
                    st.error("Dados de acesso incorretos. Por favor, tente novamente.")

# =====================================================================
# INTERFACE PRINCIPAL (MODULAR E MODERNA)
# =====================================================================
else:
    # --- BARRA LATERAL (MODERN NAV) ---
    with st.sidebar:
        # Logo da Empresa ou Ícone Padrão
        if st.session_state.logo_personalizada is not None:
            st.image(st.session_state.logo_personalizada, width=130)
        else:
            st.markdown("<h2 style='color: #1e3a8a;'>🌱 Ecoa</h2>", unsafe_allow_html=True)
        
        st.caption(f"Olá, {st.session_state.usuario_logado}")
        st.divider()
        
        # Menu de Navegação Moderno
        menu = st.radio(
            "Navegação",
            [
                "🏢 Empresa",
                "📝 Formulário da Pesquisa",
                "✉️ Mensagem Automática",
                "🔗 Gerenciamento de Links",
                "📊 Dashboard Geral",
                "📑 Relatórios",
                "👥 Clientes (Usuários)",
                "⚙️ Configurações"
            ],
            label_visibility="collapsed"
        )
        
        st.divider()
        if st.button("🚪 Encerrar Sessão", use_container_width=True):
            st.session_state.autenticado = False
            st.rerun()
            
        st.caption("v1.0.5 | Desenvolvido por Cris Lima")

    # --- TÍTULO DO MÓDULO ATUAL ---
    st.markdown(f"<h2 style='color: #334155; font-weight: 700;'>{menu}</h2>", unsafe_allow_html=True)
    st.write(f"Gestão estratégica da pesquisa de clima para: **{st.session_state.empresa_atual}**")
    st.divider()

    # =====================================================================
    # MÓDULOS (LÓGICA DETALHADA)
    # =====================================================================
    
    if menu == "🏢 Empresa":
        st.markdown("### 🏛️ Estrutura Organizacional")
        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                st.session_state.empresa_atual = st.text_input("Nome da Empresa", st.session_state.empresa_atual)
                st.text_input("CNPJ", "00.000.000/0001-00")
            with col2:
                st.selectbox("Porte Corporativo", ["Micro", "Pequena", "Média", "Grande"])
                st.selectbox("Segmento de Atuação", ["Serviços", "Indústria", "Tecnologia", "Varejo", "Saúde"])
            
            st.markdown("#### Departamentos")
            st.text_area("Mapeamento de Áreas (uma por linha)", "Recursos Humanos\nComercial\nOperações\nFinanceiro\nTI\nLogística", height=150)
            
            if st.button("💾 Salvar Informações da Empresa", type="primary"):
                st.success("Dados da empresa salvos com sucesso!")

    elif menu == "📝 Formulário da Pesquisa":
        st.markdown("### 📝 Configuração do Questionário")
        st.write("Personalize o banco de perguntas que será apresentado aos respondentes.")
        
        perg_edit = st.data_editor(
            st.session_state.perguntas,
            num_rows="dynamic",
            use_container_width=True,
            hide_index=True,
            column_config={
                "Ativa": st.column_config.CheckboxColumn("Ativa?"),
                "Pilar Estratégico": st.column_config.SelectboxColumn("Pilar", options=["Liderança & Gestão", "Segurança Psicológica", "Bem-estar & Saúde Mental", "Reconhecimento & Valorização", "Desenvolvimento & Crescimento", "Comunicação & Transparência", "Diversidade & Inclusão", "Propósito & Orgulho", "Infraestrutura & Ferramentas", "eNPS"])
            }
        )
        st.session_state.perguntas = perg_edit
        
        with st.expander("👁️ Visualizar como o Colaborador"):
            for _, p in perg_edit[perg_edit["Ativa"]].iterrows():
                if p["Pilar Estratégico"] == "eNPS":
                    st.slider(f"{p['Texto da Pergunta']}", 0, 10, 8)
                else:
                    st.select_slider(f"[{p['Pilar Estratégico']}] {p['Texto da Pergunta']}", options=["Discordo Totalmente", "Discordo", "Neutro", "Concordo", "Concordo Totalmente"], value="Concordo")

    elif menu == "✉️ Mensagem Automática":
        st.markdown("### ✉️ Convite Humanizado")
        st.write("Defina o texto de convite para a pesquisa. Use as tags dinâmicas para personalização.")
        
        body_msg = st.text_area("Texto do Convite", st.session_state.mensagem_padrao, height=200)
        st.caption("Tags: `{nome}`, `{data_validade}`, `{link_pesquisa}`")
        
        if st.button("💾 Atualizar Modelo"):
            st.session_state.mensagem_padrao = body_msg
            st.success("Modelo de mensagem salvo!")
            
        st.markdown("#### Prévia do Envio")
        prev = body_msg.replace("{nome}", "João").replace("{data_validade}", st.session_state.data_validade.strftime("%d/%m/%Y")).replace("{link_pesquisa}", "https://ecoa.app/pesquisa/demo")
        st.info(prev)

    elif menu == "🔗 Gerenciamento de Links":
        st.markdown("### 🔗 Controle de Disparos")
        c1, c2 = st.columns(2)
        with c1:
            st.session_state.data_validade = st.date_input("Data Limite da Campanha", st.session_state.data_validade)
        with c2:
            st.write(f"**Identificador da Empresa:** {st.session_state.empresa_atual.replace(' ', '').lower()}")
            
        st.divider()
        slug = st.session_state.empresa_atual.replace(" ", "").lower()
        st.markdown("#### Link Geral (Campanha)")
        st.code(f"https://ecoa.app/{slug}/pesquisa", language="html")
        
        st.markdown("#### Links Individuais")
        st.file_uploader("Subir Lista de Colaboradores (CSV/Excel)", type=["csv", "xlsx"])
        if st.button("Gerar Tokens e Disparar"):
            st.success("Disparos agendados no servidor!")

    elif menu == "📊 Dashboard Geral":
        st.markdown("### 📊 Inteligência de Dados")
        
        d_demo = st.session_state.dados_historicos
        m1, m2, m3 = st.columns(3)
        
        # Métricas em Cards
        m1.metric("Liderança (Média)", f"{d_demo['lideranca'].mean():.2f}", "+0.2")
        m2.metric("Comunicação (Média)", f"{d_demo['comunicacao'].mean():.2f}", "-0.1")
        
        # eNPS Logic
        prom = len(d_demo[d_demo['enps'] >= 9])
        detr = len(d_demo[d_demo['enps'] <= 6])
        enps_final = ((prom - detr) / len(d_demo)) * 100
        m3.metric("Score eNPS", f"{enps_final:.0f}", "Zona de Qualidade")
        
        st.divider()
        
        # Gráficos Plotly Modernos
        df_group = d_demo.groupby("departamento")[["lideranca", "comunicacao", "reconhecimento"]].mean().reset_index()
        fig = px.bar(df_group, x="departamento", y=["lideranca", "comunicacao", "reconhecimento"], 
                     barmode="group", title="Performance por Pilar e Setor",
                     color_discrete_sequence=["#1e3a8a", "#10b981", "#f59e0b"],
                     template="plotly_white")
        st.plotly_chart(fig, use_container_width=True)

    elif menu == "📑 Relatórios":
        st.markdown("### 📑 Geração de Apresentações")
        st.write("Exporte os resultados da campanha diretamente para um arquivo PowerPoint formatado.")
        
        if st.button("🚀 Criar Apresentação PPTX", type="primary"):
            prs = Presentation()
            # Capa
            s0 = prs.slides.add_slide(prs.slide_layouts[0])
            s0.shapes.title.text = f"Análise de Clima: {st.session_state.empresa_atual}"
            s0.placeholders[1].text = f"Ciclo: {datetime.now().year} | Gerado em {datetime.now().strftime('%d/%m/%Y')}"
            
            # Sumário de Médias
            s1 = prs.slides.add_slide(prs.slide_layouts[1])
            s1.shapes.title.text = "Visão Geral dos Resultados"
            tf = s1.placeholders[1].text_frame
            tf.text = "Índices Principais:"
            p1 = tf.add_paragraph(); p1.text = f"• Média Liderança: {st.session_state.dados_historicos['lideranca'].mean():.2f}"
            p2 = tf.add_paragraph(); p2.text = f"• Score eNPS: {enps_final:.0f}"
            
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            st.download_button("📥 Baixar Relatório (.pptx)", data=buffer, file_name=f"Relatorio_{st.session_state.empresa_atual}.pptx")

    elif menu == "👥 Clientes (Usuários)":
        st.markdown("### 👥 Gestão de Acessos")
        st.write("Controle os usuários administrativos e seus níveis de permissão.")
        
        users_fin = st.data_editor(
            st.session_state.usuarios,
            num_rows="dynamic",
            use_container_width=True,
            hide_index=True,
            column_config={
                "Senha": st.column_config.TextColumn("Senha", help="Senha de acesso"),
                "Perfil": st.column_config.SelectboxColumn("Perfil", options=["Admin Master", "Analista RH", "Cliente (Leitura)"]),
                "Status": st.column_config.SelectboxColumn("Status", options=["Ativo", "Inativo", "Pendente"])
            }
        )
        st.session_state.usuarios = users_fin

    elif menu == "⚙️ Configurações":
        st.markdown("### ⚙️ Preferências")
        st.markdown("#### 🎨 Customização de Logo")
        up = st.file_uploader("Substituir Imagem do Sistema", type=["png", "jpg"])
        if up:
            st.session_state.logo_personalizada = up
            st.success("Logo atualizada!")
            
        st.divider()
        st.markdown("#### 🔒 Segurança e Banco")
        if supabase:
            st.success("Conectado ao Supabase Cloud.")
        else:
            st.warning("Aguardando chaves do Supabase nos Secrets.")
        
        st.divider()
        st.caption("Ecoa - Plataforma de Escuta Organizacional | Por Cris Lima")
